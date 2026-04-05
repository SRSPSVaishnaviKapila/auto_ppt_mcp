"""
mcp_servers/pptx_server.py
══════════════════════════════════════════════════════════════════════════════
MCP SERVER 1 — PowerPoint Tools  (5 Visual Styles)
══════════════════════════════════════════════════════════════════════════════

Tools exposed:
  • create_presentation  →  Initialize a new PPT with a chosen visual style
  • add_slide            →  Add one slide in that style
  • save_presentation    →  Write the finished .pptx to disk

Available styles:
  1. "corporate"  — Navy / gold, top header band, professional
  2. "dark"       — Near-black bg, electric-cyan accents, tech/startup
  3. "minimal"    — Pure white, oversized slide numbers, editorial
  4. "gradient"   — Purple left panel + hot-pink, bold and energetic
  5. "nature"     — Forest-green / cream, organic warm feel
"""

import os
import json
import asyncio
import requests
from io import BytesIO
from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp.types import Tool, TextContent
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Module-level state ────────────────────────────────────────────────────────
_presentation = None
_slide_count  = 0
_active_style = "corporate"
_total_slides = 5
_current_slide = None  # Track the last added slide for image insertion

# ══════════════════════════════════════════════════════════════════════════════
#  Style Palette Definitions
# ══════════════════════════════════════════════════════════════════════════════

STYLES = {
    "corporate": {
        "title_bg":   (0x1E, 0x27, 0x61),
        "content_bg": (0xF4, 0xF6, 0xFF),
        "closing_bg": (0x1E, 0x27, 0x61),
        "accent":     (0xF9, 0xC7, 0x4F),
        "title_txt":  (0xFF, 0xFF, 0xFF),
        "body_txt":   (0x2D, 0x3A, 0x6B),
        "font_t":     "Calibri",
        "font_b":     "Calibri",
    },
    "dark": {
        "title_bg":   (0x0D, 0x0F, 0x1A),
        "content_bg": (0x10, 0x14, 0x22),
        "closing_bg": (0x0D, 0x0F, 0x1A),
        "accent":     (0x00, 0xE5, 0xFF),
        "title_txt":  (0xFF, 0xFF, 0xFF),
        "body_txt":   (0xCB, 0xD5, 0xE1),
        "font_t":     "Arial",
        "font_b":     "Arial",
    },
    "minimal": {
        "title_bg":   (0xFF, 0xFF, 0xFF),
        "content_bg": (0xFF, 0xFF, 0xFF),
        "closing_bg": (0xFF, 0xFF, 0xFF),
        "accent":     (0x11, 0x11, 0x11),
        "title_txt":  (0x11, 0x11, 0x11),
        "body_txt":   (0x44, 0x44, 0x44),
        "font_t":     "Georgia",
        "font_b":     "Georgia",
    },
    "gradient": {
        "title_bg":   (0x6C, 0x00, 0xBF),
        "content_bg": (0xF8, 0xF0, 0xFF),
        "closing_bg": (0x6C, 0x00, 0xBF),
        "accent":     (0xFF, 0x4D, 0xA1),
        "title_txt":  (0xFF, 0xFF, 0xFF),
        "body_txt":   (0x2D, 0x00, 0x52),
        "font_t":     "Trebuchet MS",
        "font_b":     "Trebuchet MS",
    },
    "nature": {
        "title_bg":   (0x2C, 0x5F, 0x2D),
        "content_bg": (0xF5, 0xF0, 0xE8),
        "closing_bg": (0x2C, 0x5F, 0x2D),
        "accent":     (0x97, 0xBC, 0x62),
        "title_txt":  (0xFF, 0xFF, 0xFF),
        "body_txt":   (0x3A, 0x4A, 0x30),
        "font_t":     "Cambria",
        "font_b":     "Calibri",
    },
}


# ══════════════════════════════════════════════════════════════════════════════
#  Drawing helpers
# ══════════════════════════════════════════════════════════════════════════════

def rgb(t):
    return RGBColor(*t)


def fill_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, x, y, w, h, color, line=False):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if not line:
        s.line.fill.background()
    return s


def oval(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(9, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def textbox(slide, text, x, y, w, h, size, bold, color,
            face="Calibri", align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text          = text
    run.font.size     = Pt(size)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.color.rgb = color
    run.font.name     = face
    return tb


# ══════════════════════════════════════════════════════════════════════════════
#  Style builders — one function per style
# ══════════════════════════════════════════════════════════════════════════════

def build_corporate(prs, slide, title, bullets, stype, num, total):
    """Navy header band / gold accents / clean corporate look."""
    s  = STYLES["corporate"]
    W  = prs.slide_width
    H  = prs.slide_height
    BG = rgb(s["content_bg"]) if stype == "content" else rgb(s["title_bg"])
    fill_bg(slide, BG)

    if stype in ("title", "closing"):
        # Thin gold stripe at top
        rect(slide, Inches(0), Inches(0), W, Inches(0.18), rgb(s["accent"]))
        # Thin gold stripe at bottom
        rect(slide, Inches(0), H - Inches(0.18), W, Inches(0.18), rgb(s["accent"]))
        textbox(slide, title,
                Inches(0.8), Inches(1.7), Inches(11.5), Inches(2.0),
                42, True, rgb(s["title_txt"]), s["font_t"])
        if bullets:
            sub = "  ·  ".join(bullets[:3])
            textbox(slide, sub,
                    Inches(0.8), Inches(3.8), Inches(11.0), Inches(0.8),
                    15, False, rgb(s["accent"]), s["font_b"])
    else:
        # Navy top header band
        rect(slide, Inches(0), Inches(0), W, Inches(1.2), rgb(s["title_bg"]))
        # Gold badge for slide number
        rect(slide, W - Inches(1.05), Inches(0.22),
             Inches(0.6), Inches(0.6), rgb(s["accent"]))
        textbox(slide, str(num),
                W - Inches(1.05), Inches(0.22), Inches(0.6), Inches(0.6),
                13, True, rgb(s["title_bg"]), s["font_t"], PP_ALIGN.CENTER)
        # Title in band
        textbox(slide, title,
                Inches(0.4), Inches(0.18), Inches(11.8), Inches(0.85),
                28, True, rgb(s["title_txt"]), s["font_t"])
        # Gold divider
        rect(slide, Inches(0.4), Inches(1.28), Inches(12.5), Inches(0.04),
             rgb(s["accent"]))
        # Bullet rows
        for i, b in enumerate(bullets[:5]):
            top = Inches(1.48) + i * Inches(0.82)
            oval(slide, Inches(0.35), top + Inches(0.24),
                 Inches(0.18), Inches(0.18), rgb(s["accent"]))
            textbox(slide, b, Inches(0.65), top, Inches(12.1), Inches(0.78),
                    15, False, rgb(s["body_txt"]), s["font_b"])


def build_dark(prs, slide, title, bullets, stype, num, total):
    """Near-black bg, electric-cyan accents, tech-startup style."""
    s   = STYLES["dark"]
    W   = prs.slide_width
    H   = prs.slide_height
    CYAN = rgb(s["accent"])
    WHITE = rgb(s["title_txt"])
    BODY  = rgb(s["body_txt"])

    BG = rgb(s["title_bg"]) if stype in ("title", "closing") else rgb(s["content_bg"])
    fill_bg(slide, BG)

    if stype in ("title", "closing"):
        # Cyan bottom bar
        rect(slide, Inches(0), H - Inches(0.18), W, Inches(0.18), CYAN)
        # Faint watermark number
        textbox(slide, str(num).zfill(2),
                Inches(7.5), Inches(0.2), Inches(5.5), Inches(3.8),
                120, True, RGBColor(0x18, 0x1F, 0x35),
                s["font_t"], PP_ALIGN.RIGHT)
        textbox(slide, title,
                Inches(0.7), Inches(1.8), Inches(11.0), Inches(1.8),
                40, True, WHITE, s["font_t"])
        if bullets:
            textbox(slide, bullets[0],
                    Inches(0.7), Inches(3.8), Inches(10.0), Inches(0.7),
                    16, False, CYAN, s["font_b"])
    else:
        # Cyan left stripe
        rect(slide, Inches(0), Inches(0), Inches(0.12), H, CYAN)
        textbox(slide, title,
                Inches(0.35), Inches(0.22), Inches(11.8), Inches(0.9),
                30, True, WHITE, s["font_t"])
        # Counter
        textbox(slide, f"{num}/{total}",
                Inches(11.8), Inches(0.22), Inches(1.3), Inches(0.7),
                13, False, CYAN, s["font_b"], PP_ALIGN.RIGHT)
        # Cyan divider
        rect(slide, Inches(0.35), Inches(1.18), Inches(12.3), Inches(0.03), CYAN)
        # Dark card bullets
        ch  = Inches(0.72)
        gap = Inches(0.1)
        for i, b in enumerate(bullets[:5]):
            top = Inches(1.32) + i * (ch + gap)
            rect(slide, Inches(0.35), top, Inches(12.3), ch,
                 RGBColor(0x16, 0x1C, 0x2E))
            rect(slide, Inches(0.5), top + Inches(0.28),
                 Inches(0.13), Inches(0.13), CYAN)
            textbox(slide, b, Inches(0.75), top + Inches(0.05),
                    Inches(11.6), ch, 14, False, BODY, s["font_b"])


def build_minimal(prs, slide, title, bullets, stype, num, total):
    """White bg, large numbers, black editorial typography."""
    s     = STYLES["minimal"]
    W     = prs.slide_width
    H     = prs.slide_height
    BLACK = rgb(s["accent"])
    DARK  = rgb(s["title_txt"])
    BODY  = rgb(s["body_txt"])

    fill_bg(slide, rgb(s["content_bg"]))

    if stype in ("title", "closing"):
        # Full-height black left panel
        rect(slide, Inches(0), Inches(0), Inches(4.0), H, BLACK)
        # Large faint number on black panel
        textbox(slide,
                "01" if stype == "title" else str(total).zfill(2),
                Inches(0.2), Inches(3.0), Inches(3.6), Inches(2.2),
                80, True, RGBColor(0x2A, 0x2A, 0x2A),
                s["font_t"], PP_ALIGN.RIGHT)
        textbox(slide, title,
                Inches(4.7), Inches(1.6), Inches(7.8), Inches(2.0),
                32, True, DARK, s["font_t"])
        if bullets:
            textbox(slide, bullets[0],
                    Inches(4.7), Inches(3.8), Inches(7.8), Inches(0.8),
                    14, False, RGBColor(0x66, 0x66, 0x66), s["font_b"])
    else:
        # Giant watermark number top-right
        textbox(slide, str(num).zfill(2),
                Inches(9.8), Inches(0.0), Inches(3.3), Inches(2.0),
                90, True, RGBColor(0xEE, 0xEE, 0xEE),
                s["font_t"], PP_ALIGN.RIGHT)
        # Top thin line
        rect(slide, Inches(0.5), Inches(0.45), Inches(9.0), Inches(0.03), BLACK)
        # Title
        textbox(slide, title,
                Inches(0.5), Inches(0.56), Inches(9.0), Inches(1.1),
                30, True, DARK, s["font_t"])
        # Lighter rule
        rect(slide, Inches(0.5), Inches(1.72), Inches(9.0), Inches(0.02),
             RGBColor(0xCC, 0xCC, 0xCC))
        # Numbered bullets
        for i, b in enumerate(bullets[:5]):
            top = Inches(1.92) + i * Inches(0.72)
            textbox(slide, f"{i+1}.", Inches(0.5), top, Inches(0.5), Inches(0.72),
                    15, True, BLACK, s["font_t"])
            textbox(slide, b, Inches(1.1), top, Inches(11.5), Inches(0.72),
                    14, False, BODY, s["font_b"])


def build_gradient(prs, slide, title, bullets, stype, num, total):
    """Purple left panel + hot-pink, bold high-energy."""
    s     = STYLES["gradient"]
    W     = prs.slide_width
    H     = prs.slide_height
    PURPLE = rgb(s["title_bg"])
    PINK   = rgb(s["accent"])
    WHITE  = rgb(s["title_txt"])
    BODY   = rgb(s["body_txt"])

    if stype in ("title", "closing"):
        fill_bg(slide, PURPLE)
        # Pink bars bottom-right
        rect(slide, W - Inches(4.0), H - Inches(0.38), Inches(4.0), Inches(0.38), PINK)
        rect(slide, W - Inches(6.5), H - Inches(0.18), Inches(6.5), Inches(0.18),
             RGBColor(0xA0, 0x20, 0xCC))
        textbox(slide, title,
                Inches(0.8), Inches(1.6), Inches(11.5), Inches(2.0),
                42, True, WHITE, s["font_t"])
        # Pink underline
        rect(slide, Inches(0.8), Inches(3.7), Inches(3.5), Inches(0.14), PINK)
        if bullets:
            textbox(slide, bullets[0],
                    Inches(0.8), Inches(3.98), Inches(11.0), Inches(0.8),
                    17, False, RGBColor(0xDD, 0xBB, 0xFF), s["font_b"])
    else:
        fill_bg(slide, rgb(s["content_bg"]))
        # Left purple panel ~30% width
        rect(slide, Inches(0), Inches(0), Inches(3.7), H, PURPLE)
        # Pink border on right edge of panel
        rect(slide, Inches(3.7), Inches(0), Inches(0.1), H, PINK)
        # Slide number in panel (bottom)
        textbox(slide, str(num).zfill(2),
                Inches(0.2), H - Inches(1.15), Inches(3.3), Inches(0.95),
                50, True, RGBColor(0x55, 0x00, 0x99),
                s["font_t"], PP_ALIGN.RIGHT)
        # Title on panel (vertical)
        textbox(slide, title,
                Inches(0.22), Inches(0.28), Inches(3.26), Inches(4.2),
                19, True, WHITE, s["font_t"])
        # Bullet rows on right content area
        ch  = Inches(0.80)
        gap = Inches(0.12)
        for i, b in enumerate(bullets[:5]):
            top = Inches(0.38) + i * (ch + gap)
            # Pink square marker
            rect(slide, Inches(3.98), top + Inches(0.32),
                 Inches(0.14), Inches(0.14), PINK)
            textbox(slide, b, Inches(4.22), top, Inches(8.7), ch,
                    14, False, BODY, s["font_b"])


def build_nature(prs, slide, title, bullets, stype, num, total):
    """Forest green top band / cream body, organic warm style."""
    s     = STYLES["nature"]
    W     = prs.slide_width
    H     = prs.slide_height
    GREEN = rgb(s["title_bg"])
    MOSS  = rgb(s["accent"])
    WHITE = rgb(s["title_txt"])
    BODY  = rgb(s["body_txt"])

    if stype in ("title", "closing"):
        fill_bg(slide, GREEN)
        # Lighter-green bottom arc area
        oval(slide, Inches(-1.0), H - Inches(1.6),
             Inches(15.0), Inches(3.0), RGBColor(0x22, 0x4F, 0x23))
        # Moss top stripe
        rect(slide, Inches(0), Inches(0), W, Inches(0.18), MOSS)
        textbox(slide, title,
                Inches(1.0), Inches(1.4), Inches(11.2), Inches(2.0),
                42, True, WHITE, s["font_t"], PP_ALIGN.CENTER)
        if bullets:
            sub = "  ·  ".join(bullets[:3])
            textbox(slide, f"· {sub} ·",
                    Inches(1.0), Inches(3.85), Inches(11.2), Inches(0.8),
                    14, False, MOSS, s["font_b"], PP_ALIGN.CENTER)
    else:
        fill_bg(slide, rgb(s["content_bg"]))
        # Green top band
        rect(slide, Inches(0), Inches(0), W, Inches(1.2), GREEN)
        # Leaf-dot decorations in band
        for xi in [0.28, 0.55, 0.82]:
            oval(slide, Inches(xi), Inches(0.45),
                 Inches(0.22), Inches(0.22), MOSS)
        textbox(slide, title,
                Inches(1.2), Inches(0.18), Inches(11.0), Inches(0.85),
                26, True, WHITE, s["font_t"])
        # Slide number
        textbox(slide, str(num),
                Inches(12.3), Inches(0.2), Inches(0.8), Inches(0.7),
                22, True, MOSS, s["font_t"], PP_ALIGN.RIGHT)
        # Moss divider
        rect(slide, Inches(0.5), Inches(1.3), Inches(12.0), Inches(0.05), MOSS)
        # Bullet rows with leaf dots
        ch  = Inches(0.78)
        gap = Inches(0.08)
        for i, b in enumerate(bullets[:5]):
            top = Inches(1.45) + i * (ch + gap)
            oval(slide, Inches(0.42), top + Inches(0.24),
                 Inches(0.22), Inches(0.22), MOSS)
            textbox(slide, b, Inches(0.78), top, Inches(11.8), ch,
                    15, False, BODY, s["font_b"])


BUILDERS = {
    "corporate": build_corporate,
    "dark":      build_dark,
    "minimal":   build_minimal,
    "gradient":  build_gradient,
    "nature":    build_nature,
}


# ══════════════════════════════════════════════════════════════════════════════
#  MCP Server wiring
# ══════════════════════════════════════════════════════════════════════════════

app = Server("pptx-server")


@app.list_tools()
async def list_tools():
    return [
        Tool(
            name="create_presentation",
            description="Initialize a new PowerPoint presentation with a visual style. "
                        "ALWAYS call this first.",
            inputSchema={
                "type": "object",
                "properties": {
                    "title":        {"type": "string"},
                    "style":        {
                        "type": "string",
                        "enum": list(STYLES.keys()),
                        "description": (
                            "corporate=navy/gold professional | "
                            "dark=black/cyan tech | "
                            "minimal=white editorial | "
                            "gradient=purple/pink bold | "
                            "nature=green/cream warm"
                        )
                    },
                    "total_slides": {"type": "integer"}
                },
                "required": ["title", "style", "total_slides"]
            }
        ),
        Tool(
            name="add_slide",
            description="Add one slide. Call once per slide.",
            inputSchema={
                "type": "object",
                "properties": {
                    "slide_title":   {"type": "string"},
                    "bullet_points": {"type": "array", "items": {"type": "string"}},
                    "slide_type":    {
                        "type": "string",
                        "enum": ["title", "content", "closing"]
                    }
                },
                "required": ["slide_title", "bullet_points", "slide_type"]
            }
        ),
        Tool(
            name="add_image",
            description="Add an image to the last added slide (right side or centered). "
                        "Downloads from URL and embeds in the presentation.",
            inputSchema={
                "type": "object",
                "properties": {
                    "image_url": {"type": "string", "description": "URL to the image"},
                    "width_inches": {"type": "number", "description": "Image width in inches (default 4)"},
                },
                "required": ["image_url"]
            }
        ),
        Tool(
            name="save_presentation",
            description="Save the finished presentation to disk. Call LAST.",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string"}
                },
                "required": ["filename"]
            }
        ),
    ]


@app.call_tool()
async def call_tool(name: str, arguments: dict):
    global _presentation, _slide_count, _active_style, _total_slides, _current_slide

    if name == "create_presentation":
        _presentation  = Presentation()
        _slide_count   = 0
        _active_style  = arguments.get("style", "corporate")
        _total_slides  = arguments.get("total_slides", 5)
        _presentation.slide_width  = Inches(13.33)
        _presentation.slide_height = Inches(7.5)
        return [TextContent(type="text", text=json.dumps({
            "status": "ok",
            "message": f"Created. Style='{_active_style}', topic='{arguments['title']}'"
        }))]

    elif name == "add_slide":
        if _presentation is None:
            return [TextContent(type="text", text=json.dumps(
                {"status": "error", "message": "Call create_presentation first."}))]
        _slide_count += 1
        layout = _presentation.slide_layouts[6]
        slide  = _presentation.slides.add_slide(layout)
        _current_slide = slide  # Track for image insertion
        BUILDERS.get(_active_style, build_corporate)(
            _presentation, slide,
            arguments["slide_title"],
            arguments["bullet_points"],
            arguments.get("slide_type", "content"),
            _slide_count, _total_slides
        )
        return [TextContent(type="text", text=json.dumps({
            "status": "ok",
            "message": f"Slide {_slide_count}: '{arguments['slide_title']}' done."
        }))]

    elif name == "add_image":
        if _current_slide is None:
            return [TextContent(type="text", text=json.dumps({
                "status": "error",
                "message": "No slide to add image to. Call add_slide first."
            }))]
        
        image_url = arguments.get("image_url", "")
        if not image_url:
            return [TextContent(type="text", text=json.dumps({
                "status": "error",
                "message": "image_url is required"
            }))]
        
        try:
            # Download the image from URL
            response = requests.get(image_url, timeout=10)
            if response.status_code != 200:
                return [TextContent(type="text", text=json.dumps({
                    "status": "error",
                    "message": f"Failed to download image: HTTP {response.status_code}"
                }))]
            
            # Add image to slide
            image_stream = BytesIO(response.content)
            width = Inches(arguments.get("width_inches", 4.5))
            
            # Position image on the right side of the slide
            left = Inches(13.33 - 5.0)  # Right side
            top = Inches(1.5)
            
            _current_slide.shapes.add_picture(image_stream, left, top, width=width)
            
            return [TextContent(type="text", text=json.dumps({
                "status": "ok",
                "message": f"Image added to slide {_slide_count}."
            }))]
        
        except requests.exceptions.RequestException as e:
            return [TextContent(type="text", text=json.dumps({
                "status": "error",
                "message": f"Failed to download image: {str(e)}"
            }))]
        except Exception as e:
            return [TextContent(type="text", text=json.dumps({
                "status": "error",
                "message": f"Failed to add image to slide: {str(e)}"
            }))]

    elif name == "save_presentation":
        if _presentation is None:
            return [TextContent(type="text", text=json.dumps(
                {"status": "error", "message": "No presentation to save."}))]
        fn = arguments["filename"]
        if not fn.endswith(".pptx"):
            fn += ".pptx"
        out = os.path.join(
            os.path.dirname(os.path.dirname(__file__)), "outputs"
        )
        os.makedirs(out, exist_ok=True)
        fp = os.path.join(out, fn)
        _presentation.save(fp)
        return [TextContent(type="text", text=json.dumps({
            "status": "ok", "message": "Saved!", "filepath": fp,
            "slides": _slide_count, "style": _active_style
        }))]

    return [TextContent(type="text", text=json.dumps(
        {"status": "error", "message": f"Unknown tool: {name}"}))]


async def main():
    async with stdio_server() as (r, w):
        await app.run(r, w, app.create_initialization_options())


if __name__ == "__main__":
    asyncio.run(main())
