"""
agent/executor.py
══════════════════════════════════════════════════════════════════════════════
AGENT STEP 2 — Slide Executor (Agentic Loop)
══════════════════════════════════════════════════════════════════════════════

For each slide in the plan:
  1. (Optional) Call search MCP server to get Wikipedia facts
  2. (Optional) Ask LLM to enrich the bullet points with those facts
  3. Call PPTX MCP server to add the styled slide
  4. (Optional) Fetch a relevant image and add it
After the loop: save the file.
"""

import json
import re
import os
import asyncio
from mcp import ClientSession, StdioServerParameters
from mcp.client.stdio import stdio_client
from agent.llm_client import ask_llm


ENRICH_PROMPT = """You are a slide content writer.

PRESENTATION TOPIC: {topic}
SLIDE TITLE: "{slide_title}"
AUDIENCE: {audience}

SEARCH RESULTS (use relevant facts; ignore if unhelpful):
{search_content}

ORIGINAL BULLET POINTS:
{original_bullets}

TASK: Rewrite or improve the bullet points using the search results.
Keep them focused on the presentation topic: "{topic}"
Keep 3–5 bullets. Each must be one clear, factual, complete sentence.
Do NOT include bullet symbols. Return ONLY a JSON array of strings.

Example: ["Sentence one.", "Sentence two.", "Sentence three."]

JSON array only:
"""


def _extract_topic(user_request: str) -> str:
    """Extract the main topic/keyword from the user prompt."""
    # Remove common words and get the core topic
    words = user_request.lower().split()
    # Filter out common stop words and very short words
    stop_words = {"a", "an", "the", "for", "about", "tell", "me", "explain", 
                  "on", "in", "of", "and", "or", "is", "are", "be", "create",
                  "make", "write", "generate", "build", "show", "display", "add",
                  "list", "give", "provide", "describe"}
    
    topic_words = [w.strip('.,!?;:') for w in words 
                   if w.strip('.,!?;:') not in stop_words and len(w) > 2]
    
    if topic_words:
        # Return first few meaningful words as topic
        return " ".join(topic_words[:4])
    return user_request[:50]


def _detect_audience(user_request: str) -> str:
    """Extract audience hint from the user prompt."""
    lower = user_request.lower()
    for keyword in ["grade", "elementary", "middle school", "high school",
                    "college", "university", "professional", "beginner", "expert"]:
        if keyword in lower:
            return f"mentioned audience: '{keyword}'"
    return "general audience"


async def execute_plan(
    slide_plan:       list[dict],
    style:            str,
    user_request:     str,
    output_filename:  str,
    use_search:       bool = True,
    log_callback=None
) -> str:
    """
    Execute the slide plan via MCP tool calls.

    Parameters
    ----------
    slide_plan      : From planner.py
    style           : Visual style name, e.g. "dark"
    user_request    : Original user prompt
    output_filename : Output .pptx filename (no path)
    use_search      : Whether to call the search MCP server
    log_callback    : callable(str) for live UI logs

    Returns
    -------
    str : Absolute path to saved .pptx, or "" on failure
    """
    def log(msg):
        if log_callback:
            log_callback(msg)
    
    # Extract the main topic for consistent context throughout execution
    main_topic = _extract_topic(user_request)

    audience = _detect_audience(user_request)

    log("🔌 **Step 2: Connecting to MCP servers...**")

    base_dir      = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    pptx_script   = os.path.join(base_dir, "mcp_servers", "pptx_server.py")
    search_script = os.path.join(base_dir, "mcp_servers", "search_server.py")

    pptx_params   = StdioServerParameters(command="python", args=[pptx_script])
    search_params = StdioServerParameters(command="python", args=[search_script])

    try:
        async with stdio_client(pptx_params) as (pr, pw), \
                   stdio_client(search_params) as (sr, sw):
            async with ClientSession(pr, pw) as pptx, \
                       ClientSession(sr, sw) as search:

                await pptx.initialize()
                await search.initialize()
                log("✅ Both MCP servers connected.")

                # ── create_presentation ────────────────────────────────────────
                log(f"\n📂 Creating presentation | Style: **{style}**")
                res = await pptx.call_tool("create_presentation", {
                    "title":        user_request,
                    "style":        style,
                    "total_slides": len(slide_plan)
                })
                log(f"   ↳ {_msg(res)}")

                # ── Agentic slide loop ────────────────────────────────────────
                log(f"\n🔄 **Building {len(slide_plan)} slides...**")

                for slide in slide_plan:
                    num   = slide["slide_number"]
                    ttl   = slide["slide_title"]
                    stype = slide["slide_type"]
                    bullets = list(slide["bullet_points"])

                    log(f"\n─── Slide {num}/{len(slide_plan)}: \"{ttl}\" [{stype}] ───")

                    # Optional: enrich content slides with Wikipedia facts
                    if use_search and stype == "content":
                        log(f"   🔍 Searching Wikipedia: '{ttl}'")
                        sr_res = await search.call_tool("search_web", {"query": ttl})
                        data   = _parse(sr_res)
                        content = data.get("content", "")

                        if content:
                            log(f"   📖 Got {len(content)} chars — enriching bullets...")
                            bullets = _enrich(ttl, bullets, content, audience, main_topic)
                            log("   ✨ Bullets enriched.")
                        else:
                            log("   ℹ️  No Wikipedia result — using LLM knowledge.")
                    else:
                        log(f"   ⏭️  Skipping search (type: {stype})")

                    # add_slide
                    add_res = await pptx.call_tool("add_slide", {
                        "slide_title":   ttl,
                        "bullet_points": bullets,
                        "slide_type":    stype
                    })
                    log(f"   🖊️  {_msg(add_res)}")

                    # Fetch and add image (for content and closing slides)
                    if stype in ("content", "closing"):
                        # Combine slide title with main topic for better image search
                        image_keyword = f"{main_topic} {ttl}" if main_topic else ttl
                        log(f"   🖼️  Fetching image for: '{image_keyword}'")
                        img_res = await search.call_tool("fetch_image", {"keyword": image_keyword})
                        img_data = _parse(img_res)
                        img_url = img_data.get("image_url", "")
                        
                        if img_url:
                            log(f"   📸 Got image — adding to slide...")
                            add_img_res = await pptx.call_tool("add_image", {
                                "image_url": img_url,
                                "width_inches": 4.5
                            })
                            log(f"   ✅ {_msg(add_img_res)}")
                        else:
                            log(f"   ℹ️  No image found for this slide.")

                # ── save_presentation ─────────────────────────────────────────
                log("\n💾 **Saving...**")
                save_res = await pptx.call_tool("save_presentation",
                                                {"filename": output_filename})
                save_data = _parse(save_res)

                if save_data.get("status") == "ok":
                    fp = save_data.get("filepath", "")
                    log(f"\n🎉 Done! {save_data.get('slides')} slides saved to:\n   `{fp}`")
                    return fp
                else:
                    log(f"❌ Save failed: {save_data}")
                    return ""

    except Exception as e:
        log(f"❌ Executor error: {e}")
        log("🔄 Attempting fallback (direct python-pptx)...")
        return _fallback(slide_plan, output_filename, log)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _enrich(title, original, search_content, audience, topic):
    prompt = ENRICH_PROMPT.format(
        topic=topic,
        slide_title=title,
        audience=audience,
        search_content=search_content[:400],
        original_bullets="\n".join(f"- {b}" for b in original)
    )
    raw = ask_llm(prompt, temperature=0.4)
    m   = re.search(r'\[.*?\]', raw, re.DOTALL)
    if m:
        try:
            bullets = json.loads(m.group())
            if isinstance(bullets, list) and all(isinstance(b, str) for b in bullets):
                return [b.strip() for b in bullets[:5] if b.strip()] or original
        except json.JSONDecodeError:
            pass
    return original


def _parse(result) -> dict:
    try:
        content = result.content
        if content and hasattr(content[0], "text"):
            return json.loads(content[0].text)
    except Exception:
        pass
    return {}


def _msg(result) -> str:
    d = _parse(result)
    return d.get("message", d.get("status", str(d)))


def _fallback(slide_plan, output_filename, log) -> str:
    """Direct python-pptx fallback if MCP servers fail to start."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        for sd in slide_plan:
            layout = prs.slide_layouts[1]
            slide  = prs.slides.add_slide(layout)
            slide.shapes.title.text = sd["slide_title"]
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, b in enumerate(sd["bullet_points"]):
                if i == 0:
                    tf.paragraphs[0].text = b
                else:
                    tf.add_paragraph().text = b
        base = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        out  = os.path.join(base, "outputs")
        os.makedirs(out, exist_ok=True)
        fn = output_filename if output_filename.endswith(".pptx") else output_filename + ".pptx"
        fp = os.path.join(out, fn)
        prs.save(fp)
        log(f"✅ Fallback saved to: {fp}")
        return fp
    except Exception as e:
        log(f"❌ Fallback failed: {e}")
        return ""
